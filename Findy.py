import customtkinter
import tkinter
from tkinter.filedialog import askdirectory
import os
import mimetypes
import docx2txt
import pdfplumber
from pptx import Presentation
import openpyxl as xl
from pyxlsb import open_workbook
import xlrd
import win32com.client
import docx2txt
"""
This app was made by Baje99. 
please assume that the application is currently a demo, and for any error you can contact me on github to be able to fix it.
Supported Formats:
Word: '.doc', '.docx'
Text: '.txt'
PDF: '.pdf'
Powerpoint: '.ppt', '.pptx'
Excel: '.xlsx', '.csv', '.xlsm', '.xltx', '.xltm', '.xls', '.xlsb'

FindyApp will walk through selected directory and subdirectories, filter the files based on selected extension, open the files and check if the word is found.
Being still untested in the long term, it is possible to be susceptible to bugs and not to take data from files containing complex elements.

It is preferable to be used for simple files that do not contain tables, graphs, functions, etc"""

def findWord(filetype, word):
    wordsfound = []
    if filetype == "Word/Pdf/Txt":
        filetypelist = ['.doc', '.docx', '.txt', '.pdf']
    elif filetype == 'Powerpoint':
        filetypelist = ['.pptx', '.ppt']
    elif filetype == "Excel":
        filetypelist = ['.xlsx', '.csv', '.xlsm', '.xltx', '.xltm', '.xls', '.xlsb']
    else:
        filetypelist = ['.pptx', '.ppt', '.doc', '.docx', '.txt', '.pdf', '.xlsx', '.csv', '.xlsm', '.xltx', '.xltm', '.xls', '.xlsb']
    for dir, subdirs, files in os.walk(path1):
        for filename in files:
            filepath = os.path.join(dir, filename)
            extension = os.path.splitext(filepath)[1]
            if not extension:
                extension = mimetypes.guess_extension(filepath)
            if extension in filetypelist:
                print(filepath)
                if filepath.lower().find(word) >= 0:
                        print("{} was found in the {} file".format(word,filepath))
                        wordsfound.append(filepath)
                elif extension == ".docx":
                    content = docx2txt.process(filepath).lower()
                    if content.find(word)>=0:
                        print("{} was found in the {} file".format(word,filepath))
                        wordsfound.append(filepath)
                elif extension == ".doc":
                    wordclient = win32com.client.Dispatch("Word.application")
                    wordDoc = wordclient.Documents.Open(filepath, False, False, False)
                    newfilepath = filepath.replace(".doc", ".docx")
                    wordDoc.SaveAs2(newfilepath, FileFormat = 16)
                    wordDoc.Close()
                    content = docx2txt.process(newfilepath).lower()
                    if content.find(word) >=0 :
                        print("{} was found in the {} file".format(word,filepath))
                        wordsfound.append(filepath)
                    try: 
                        os.remove(newfilepath)
                    except:
                        pass
                elif extension == ".pdf":
                    with pdfplumber.open(filepath) as pdf:
                        for x in range(0, len(pdf.pages)):
                            page = pdf.pages[x]
                            if page.extract_text().lower().find(word)>=0:
                                wordsfound.append(filepath)
                                print("{} was found in the {} file".format(word,filepath))
                                break
                elif extension == ".txt":
                    flag = False
                    file = open(filepath, "r")
                    for line in file.readlines():
                        if line.lower().find(word)>=0:
                            wordsfound.append(filepath)
                            print("{} was found in the {} file".format(word,filepath))
                            flag = True
                            break
                    if flag: break
                elif extension == ".ppt" or extension == ".pptx":
                    flag = False
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                if shape.text.lower().find(word)>=0:
                                    wordsfound.append(filepath)
                                    print("{} was found in the {} file".format(word,filepath))
                                    flag = True
                                    break
                        if flag: break
                elif extension == ".xlsx" or extension == ".xlsm" or extension == ".xltx" or extension == ".xltm":
                    flag = False
                    workbook = xl.load_workbook(filename=filepath)
                    wsname = workbook.sheetnames[0]
                    ws = workbook[wsname]
                    for row in ws.iter_rows( max_row = ws.max_row, max_col = ws.max_column, values_only = True):
                        for value in row:
                            if str(value).lower().find(word) >= 0:
                                wordsfound.append(filepath)
                                print("{} was found in the {} file".format(word,filepath))
                                flag = True
                                break
                        if flag: break
                elif extension == ".xls":
                    flag = False
                    workbook = xlrd.open_workbook(filepath)
                    ws = workbook.sheet_by_index(0)
                    for i in range(ws.nrows):
                        for j in range(ws.ncols):
                            if str(ws.cell_value(i,j)).lower().find(word) >= 0:
                                wordsfound.append(filepath)
                                print("{} was found in the {} file".format(word,filepath))
                                flag = True
                                break
                        if flag: break
                elif extension == ".xlsb":
                    flag = False
                    with open_workbook(filepath) as wb:
                        ws = wb.get_sheet(1)
                        for row in ws.rows():
                            for col in row:
                                colstr = str(col).replace(')', '')
                                descr = colstr.split('v=')[1]
                                if descr.lower().find(word) >=0:
                                    wordsfound.append(filepath)
                                    print("{} was found in the {} file".format(word,filepath))
                                    flag = True
                                    break
                            if flag: break            
    for item in wordsfound:
        print(item)
    
def ChooseDir():
    global path1
    path1 = askdirectory(mustexist = True)
    print(path1)        

def CreateInterface():
    customtkinter.set_appearance_mode("dark")  # Modes: system (default), light, dark
    customtkinter.set_default_color_theme("green")  # Themes: blue (default), dark-blue, green
    app = customtkinter.CTk()
    app.geometry("520x400")
    app.title("FindyApp")
    app.iconbitmap(r"C:\Users\bajan\Desktop\Informatica\Python\MarketingBot\terminator.ico")
    frame_1 = customtkinter.CTkFrame(master=app)
    frame_1.pack(pady=20, padx=60, fill="both", expand=True)

    label1 = customtkinter.CTkLabel(master=frame_1, text="File type:", justify=tkinter.LEFT)
    label1.pack(side = tkinter.TOP)
    FileType = ["Any", 'Word/Pdf/Txt', "Powerpoint", "Excel"]
    entryvalue1 = tkinter.StringVar()
    entry1 = customtkinter.CTkComboBox(master = frame_1, values = FileType, variable = entryvalue1)
    entry1['state'] = 'readonly'
    entry1.set("Choose File Type:")
    entry1.pack(pady = 6, padx = 10)

    label2 = customtkinter.CTkLabel(master=frame_1, text="Enter the word to search:", justify=tkinter.LEFT)
    entryvalue2 = tkinter.StringVar()
    entry2 = customtkinter.CTkEntry(master=frame_1, textvariable=entryvalue2, placeholder_text="Enter a string value:")
    label2.pack(side = tkinter.TOP)
    entry2.pack(pady=6, padx=10)

    label3 = customtkinter.CTkLabel(master=frame_1, text='ChooseTheDirectory: ', justify=tkinter.LEFT)
    browse3 = customtkinter.CTkButton(master=frame_1, text='Browse Folder', width=10,command = lambda:ChooseDir())
    label3.pack(side = tkinter.TOP)
    browse3.pack(pady=6, padx=10)

    button1 = customtkinter.CTkButton(master=frame_1, text="Submit", command = lambda:findWord(entryvalue1.get(), entryvalue2.get().lower()))
    button2 = customtkinter.CTkButton(master=frame_1, text='Quit', command=app.quit)
    button1.pack(pady=6, padx=10)
    button2.pack(pady=6, padx=10)

    label4 = customtkinter.CTkLabel(master=frame_1, text='Github: Baje99', justify=tkinter.LEFT)
    label4.pack(side = tkinter.BOTTOM)
    app.mainloop()

if __name__ == "__main__":
    CreateInterface()